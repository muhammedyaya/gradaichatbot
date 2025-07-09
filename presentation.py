from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os
import re
from typing import List, Dict, Optional, Union

def is_arabic(text: str) -> bool:
    return any('\u0600' <= char <= '\u06FF' for char in text)

def add_formatted_paragraph(text_frame, text, font_name, font_size, alignment, color_rgb):
    p = text_frame.add_paragraph()
    p.alignment = alignment
    parts = re.split(r"(\*\*.*?\*\*)", text)
    is_rtl = is_arabic(text)
    p.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
    p.rtl = is_rtl

    for part in parts:
        run = p.add_run()
        clean = part.strip("*")
        run.text = clean
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color_rgb
        run.font.bold = part.startswith("**") and part.endswith("**")

def create_presentation(
    slides: List[Dict],
    language: str = "arabic",
    template_path: str = "template.pptx",
    theme_color: str = "#2A5CAA",  # optional fallback
    max_bullets: int = 5,
    output_format: str = "bytes",
    title_font_size: int = 32,
    bullet_font_size: int = 22,
    custom_fonts: Optional[Dict] = None,
    title_color: str = "#2A5CAA",
    bullet_color: str = "#2A5CAA"
) -> Union[BytesIO, bytes]:

    if not slides:
        raise ValueError("Slides list cannot be empty")

    for i, slide in enumerate(slides, 1):
        if not isinstance(slide, dict):
            raise ValueError(f"Slide {i} must be a dictionary")
        if "title" not in slide or "bullets" not in slide:
            raise ValueError(f"Slide {i} missing 'title' or 'bullets'")

    fonts = {"arabic": "Arial", "english": "Calibri"}
    if custom_fonts:
        fonts.update(custom_fonts)

    try:
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")

        prs = Presentation(template_path)

        title_rgb = RGBColor.from_string(title_color.lstrip("#"))
        bullet_rgb = RGBColor.from_string(bullet_color.lstrip("#"))
        is_rtl = language.lower() == "arabic"
        alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
        font_name = fonts.get(language.lower(), fonts["english"])

        TITLE_WIDTH = Inches(8)
        CONTENT_WIDTH = Inches(7.5)
        LEFT_MARGIN = Inches(1.8) if is_rtl else Inches(1)

        blank_layout = next(
            (layout for layout in prs.slide_layouts if "blank" in layout.name.lower()),
            prs.slide_layouts[0]
        )

        for i, slide_data in enumerate(slides):
            slide = prs.slides.add_slide(blank_layout)

            # Title
            title_shape = slide.shapes.add_textbox(
                LEFT_MARGIN, Inches(0.5), TITLE_WIDTH, Inches(1.5)
            )
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True

            add_formatted_paragraph(
                title_frame,
                slide_data.get("title", f"Slide {i+1}"),
                font_name,
                title_font_size,
                alignment,
                title_rgb
            )

            # Bullets
            bullets = slide_data.get("bullets", [])[:max_bullets]
            if bullets:
                content_shape = slide.shapes.add_textbox(
                    LEFT_MARGIN, Inches(1.6), CONTENT_WIDTH, Inches(4)
                )
                content_frame = content_shape.text_frame
                content_frame.word_wrap = True

                for bullet in bullets:
                    add_formatted_paragraph(
                        content_frame,
                        bullet,
                        font_name,
                        bullet_font_size,
                        alignment,
                        bullet_rgb
                    )

            # Slide number
            if i > 0:
                num_box = slide.shapes.add_textbox(
                    prs.slide_width - Inches(1),
                    prs.slide_height - Inches(0.5),
                    Inches(1),
                    Inches(0.4)
                )
                tf = num_box.text_frame
                tf.text = str(i + 1)
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Cleanup blank slides
        while len(prs.slides) > len(slides):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        output = BytesIO()
        prs.save(output)

        if output_format == "bytes":
            return output.getvalue()
        return output

    except Exception as e:
        raise RuntimeError(f"Presentation creation failed: {str(e)}")
    finally:
        if 'output' in locals():
            output.close()
